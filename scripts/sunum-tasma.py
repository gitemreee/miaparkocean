"""Metin kutusu taşma denetimi — sunumdaki her kutuyu ÖLÇER.

Render'a bakmak yetmiyor: LibreOffice'in satır kırması ile PowerPoint'inki
birebir aynı değil. Burada her kutu metrik olarak ölçülüp kutuya sığıp
sığmadığı hesaplanıyor.

Yazı tipi Calibri; bu makinede Carlito kurulu (Calibri ile metrik uyumlu,
Google'ın açık kaynak eşleniği). Carlito bulunamazsa Liberation Sans'a
düşülüyor — o Calibri'den GENİŞ, yani ölçüm güvenli tarafta hata yapar.

    python3 scripts/sunum-tasma.py sunum/MIA-PARK-OCEAN-Lansman-Sunumu.pptx
"""
import os
import sys
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
# Sunumun gerçek yazı tipleri; ölçüm bunlarla yapılıyor.
YAZI = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                    "sunum", "yazitipi")
AILE = {
    "playfair display": ("PlayfairDisplay-Regular.ttf", "PlayfairDisplay-Bold.ttf"),
    "montserrat":       ("Montserrat-Regular.ttf",      "Montserrat-Bold.ttf"),
}
LIB = "/usr/share/fonts/truetype/liberation/LiberationSans-%s.ttf"
PX = 4  # punto başına piksel (ölçek; oranlar korunur)


def font(face, size_pt, bold):
    ad = (face or "").strip().lower()
    dosyalar = AILE.get(ad)
    yol = ""
    if dosyalar:
        yol = os.path.join(YAZI, dosyalar[1 if bold else 0])
    if not os.path.exists(yol):
        yol = LIB % ("Bold" if bold else "Regular")   # geniş — güvenli taraf
    return ImageFont.truetype(yol, max(4, round(size_pt * PX)))


def wrap(words, f, max_px, draw):
    lines, cur = 1, ""
    for w in words:
        t = (cur + " " + w).strip()
        if f.getlength(t) <= max_px or not cur:
            cur = t
        else:
            lines += 1
            cur = w
    return lines


def main(path):
    prs = Presentation(path)
    n_slayt = len(prs.slides._sldIdLst)
    bad = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            bw = sh.width / EMU * 72 * PX
            bh = sh.height / EMU * 72 * PX
            tf = sh.text_frame
            ml = (tf.margin_left + tf.margin_right) / EMU * 72 * PX
            mt = (tf.margin_top + tf.margin_bottom) / EMU * 72 * PX
            avail_w, avail_h = bw - ml, bh - mt
            total_h, widest = 0.0, 0.0
            for para in tf.paragraphs:
                runs = [r for r in para.runs if r.text]
                if not runs:
                    continue
                size = max((r.font.size.pt for r in runs if r.font.size), default=18)
                bold = any(r.font.bold for r in runs)
                face = next((r.font.name for r in runs if r.font.name), None)
                f = font(face, size, bold)
                txt = "".join(r.text for r in runs)
                n = wrap(txt.split(), f, avail_w, None)
                widest = max(widest, min(f.getlength(txt), avail_w))
                total_h += n * size * PX * 1.22
            if total_h > avail_h + 1:
                bad.append((i, sh.text_frame.text.strip()[:44].replace("\n", " / "),
                            round(total_h / PX / 72, 2), round(avail_h / PX / 72, 2)))
    if not bad:
        print("TAŞMA YOK — %d slaytta bütün metin kutuları sığıyor." % n_slayt)
        return 0
    print(f"{len(bad)} kutu taşıyor:")
    for s, t, need, have in bad:
        print(f"  slayt {s:>2}  gereken {need}\"  kutu {have}\"   « {t} »")
    return 1


sys.exit(main(sys.argv[1]))
