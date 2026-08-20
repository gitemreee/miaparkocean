#!/usr/bin/env python3
"""
Çakışma denetimi — kutusuna sığan ama KOMŞUSUNUN üstüne binen metinleri bulur.

sunum-tasma.py bir metnin kendi kutusuna sığıp sığmadığına bakar. Ama asıl
bozan hata başka: iki satıra sarınca kutusundan taşmayıp bir alttaki altın
çizginin ya da bir sonraki metnin üstüne oturan başlıklar. Burada her metin
çerçevesinin GERÇEK yüksekliği (satır sayısı × satır yüksekliği) hesaplanıp
o dikdörtgen komşularıyla kesiştiriliyor.

Filled dikdörtgenlerle kesişme normaldir (blok zeminler), rapor edilmez.
Metin ↔ metin ve metin ↔ çizgi kesişmeleri rapor edilir.

    python3 scripts/sunum-cakisma.py sunum/MIA-PARK-OCEAN-Emlakci-Sunumu.pptx
"""

import os
import sys
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
PX = 4                      # punto başına piksel (ölçek; oranlar korunur)
KOK = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
YAZI = os.path.join(KOK, "sunum", "yazitipi")
AILE = {
    "playfair display": ("PlayfairDisplay-Regular.ttf", "PlayfairDisplay-Bold.ttf"),
    "montserrat":       ("Montserrat-Regular.ttf",      "Montserrat-Bold.ttf"),
}
LIB = "/usr/share/fonts/truetype/liberation/LiberationSans-%s.ttf"
TOLERANS = 0.03             # inç — bu kadar örtüşme göz ardı edilir


def font(face, boy, kalin):
    d = AILE.get((face or "").strip().lower())
    yol = os.path.join(YAZI, d[1 if kalin else 0]) if d else ""
    if not os.path.exists(yol):
        yol = LIB % ("Bold" if kalin else "Regular")
    return ImageFont.truetype(yol, max(4, round(boy * PX)))


def satir_sayisi(metin, f, genislik_px):
    """Metnin verilen genişlikte kaç satıra sarındığı."""
    n = 0
    for parca in metin.split("\n"):
        kelimeler = parca.split()
        if not kelimeler:
            n += 1
            continue
        n += 1
        cur = ""
        for k in kelimeler:
            t = (cur + " " + k).strip()
            if f.getlength(t) <= genislik_px or not cur:
                cur = t
            else:
                n += 1
                cur = k
    return n


def cerceveler(slayt):
    """(tur, x, y, w, h, etiket) listesi. tur: metin | cizgi | blok"""
    out = []
    for sh in slayt.shapes:
        try:
            x, y = sh.left / EMU, sh.top / EMU
            w, h = sh.width / EMU, sh.height / EMU
        except TypeError:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame
            metin = tf.text
            boy, kalin, face = 12.0, False, "Montserrat"
            for par in tf.paragraphs:
                for r in par.runs:
                    if r.font.size:
                        boy = r.font.size.pt
                    if r.font.bold is not None:
                        kalin = r.font.bold
                    if r.font.name:
                        face = r.font.name
                    break
                break
            f = font(face, boy, kalin)
            n = satir_sayisi(metin, f, max(1, w * 72 * PX))
            # satır yüksekliği: punto × 1.20 (PowerPoint tek satır aralığı)
            gercek = n * boy * 1.20 / 72.0
            # dolgusu olan kutular zemin bloğu olabilir; metin kutusu say
            out.append(("metin", x, y, w, max(gercek, 0.10),
                        metin.replace("\n", " ")[:44]))
        else:
            # çizgi mi blok mu: yüksekliği ya da genişliği ~0 olan çizgidir
            tur = "cizgi" if (h < 0.02 or w < 0.02) else "blok"
            out.append((tur, x, y, max(w, 0.01), max(h, 0.01), ""))
    return out


def kesisim(a, b):
    ax0, ay0, ax1, ay1 = a[1], a[2], a[1] + a[3], a[2] + a[4]
    bx0, by0, bx1, by1 = b[1], b[2], b[1] + b[3], b[2] + b[4]
    dx = min(ax1, bx1) - max(ax0, bx0)
    dy = min(ay1, by1) - max(ay0, by0)
    return dx, dy


def main(yol):
    prs = Presentation(yol)
    bulgu = []
    for i, sl in enumerate(prs.slides, 1):
        cs = cerceveler(sl)
        metinler = [c for c in cs if c[0] == "metin"]
        cizgiler = [c for c in cs if c[0] == "cizgi"]
        for a in range(len(metinler)):
            for b in range(a + 1, len(metinler)):
                dx, dy = kesisim(metinler[a], metinler[b])
                if dx > TOLERANS and dy > TOLERANS:
                    bulgu.append((i, "metin↔metin", round(dy, 2),
                                  metinler[a][5], metinler[b][5]))
        for m in metinler:
            for c in cizgiler:
                dx, dy = kesisim(m, c)
                if dx > TOLERANS and dy > 0.005:
                    bulgu.append((i, "metin↔çizgi", round(dy, 3), m[5], "—"))
        # slayt dışına taşma
        for m in metinler:
            if m[2] + m[4] > 7.5 + TOLERANS:
                bulgu.append((i, "slayt dışı", round(m[2] + m[4] - 7.5, 2), m[5], "—"))
            if m[1] + m[3] > 13.333 + TOLERANS:
                bulgu.append((i, "sağdan taşma", round(m[1] + m[3] - 13.333, 2), m[5], "—"))

    if not bulgu:
        print("ÇAKIŞMA YOK — bütün metin çerçeveleri temiz.")
        return 0
    print("%d çakışma:" % len(bulgu))
    for b in bulgu:
        print("  slayt %2d  %-12s %5.2f\"  « %s »%s"
              % (b[0], b[1], b[2], b[3], ("  ↔  « %s »" % b[4]) if b[4] != "—" else ""))
    return 1


if __name__ == "__main__":
    sys.exit(main(sys.argv[1]))
